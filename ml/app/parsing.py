"""Document parsing: PDF / DOCX / PPTX bytes -> (page_number, text) pairs."""
from __future__ import annotations

import io


def parse(filename: str, data: bytes) -> list[tuple[int, str]]:
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    if ext not in ("pdf", "docx", "pptx", "ppt"):
        raise ValueError(f"Định dạng tệp không hỗ trợ: '.{ext}'.")

    try:
        if ext == "pdf":
            return _parse_pdf(data)
        if ext == "docx":
            return _parse_docx(data)
        return _parse_pptx(data)
    except Exception as exc:  # noqa: BLE001 - corrupt/invalid file is a data error (400), not a server error
        raise ValueError(
            f"Không đọc được tệp .{ext}: tệp có thể bị hỏng hoặc sai định dạng ({type(exc).__name__})."
        ) from exc


def _parse_pdf(data: bytes) -> list[tuple[int, str]]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    return [(i + 1, (page.extract_text() or "").strip()) for i, page in enumerate(reader.pages)]


def _parse_docx(data: bytes) -> list[tuple[int, str]]:
    import docx

    document = docx.Document(io.BytesIO(data))
    text = "\n".join(p.text for p in document.paragraphs if p.text.strip())
    return [(1, text)]


def _parse_pptx(data: bytes) -> list[tuple[int, str]]:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    pages: list[tuple[int, str]] = []
    for index, slide in enumerate(presentation.slides):
        parts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        pages.append((index + 1, "\n".join(parts)))
    return pages
